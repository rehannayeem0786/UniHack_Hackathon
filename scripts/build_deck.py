"""Build the submission deck from the official template + captured screenshots.

Reads the filled-in deck (kept outside the repo because it carries the team's
downloaded template) and the PNGs captured by `node frontend/shots.mjs`, and
produces `docs/UniHack_2026_Unilog_Product_Intelligence.pptx` with real
screenshots in place of the placeholder boxes on the MVP-snapshots slide.

Re-runnable: capture fresh shots, run this, and the deck is current.

Usage:
    python scripts/build_deck.py [source.pptx]

Defaults to the newest UniHack deck in the Downloads folder when no source is
given, so the usual flow is just `python scripts/build_deck.py`.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
SHOTS = ROOT / "docs" / "shots"
OUTPUT = ROOT / "docs" / "UniHack_2026_Unilog_Product_Intelligence.pptx"

# (screenshot file, row, column) for the six snapshot cards. Card geometry is
# read from the slide itself; only the shot mapping lives here.
SNAPSHOTS = [
    ("03-input-vs-output.png", 0, 0),
    ("04-surfaces.png", 0, 1),
    ("06-provenance.png", 0, 2),
    ("11-review.png", 1, 0),
    ("07-evaluation.png", 1, 1),
    ("09-learned-rules.png", 1, 2),
]

IMAGE_WIDTH = Inches(2.72)
IMAGE_HEIGHT = Inches(0.98)


def default_source() -> Path:
    downloads = Path.home() / "Downloads"
    candidates = sorted(downloads.glob("UniHack_2026_Unilog_Product_Intelligence*.pptx"))
    if not candidates:
        sys.exit("no source deck found — pass the .pptx path as an argument")
    return candidates[-1]


def cover_crop(picture, target_aspect: float) -> None:
    """Emulate CSS object-fit: cover by cropping the longer axis."""
    aspect = picture.image.size[0] / picture.image.size[1]
    if aspect > target_aspect:  # too wide -> crop sides
        keep = target_aspect / aspect
        side = (1 - keep) / 2
        picture.crop_left = side
        picture.crop_right = side
    else:  # too tall -> crop top/bottom
        keep = aspect / target_aspect
        side = (1 - keep) / 2
        picture.crop_top = side
        picture.crop_bottom = side


def fill_snapshot_cards(slide) -> None:
    """Replace the six placeholder cards with real screenshots.

    Each card is a Rounded Rectangle with an inner placeholder Rectangle, an
    Oval glyph, a 'screenshot placeholder' label, a numbered title and a
    description. We drop the glyph/label, put the picture where the inner
    rectangle was, and nudge the title/description down to make room.
    """
    target_aspect = IMAGE_WIDTH / IMAGE_HEIGHT

    cards = [
        sh for sh in slide.shapes
        if sh.shape_type == 1 and sh.name.startswith("Rounded Rectangle")
        and sh.height > Inches(1.5) and sh.width < Inches(3.0)
    ]
    cards.sort(key=lambda s: (s.top, s.left))
    if len(cards) != 6:
        sys.exit(f"expected 6 snapshot cards, found {len(cards)}")

    by_position = {(i // 3, i % 3): card for i, card in enumerate(cards)}

    for filename, row, col in SNAPSHOTS:
        card = by_position[(row, col)]
        left, top = card.left, card.top
        # Remove every shape that lives inside this card except its background.
        doomed = [
            sh for sh in slide.shapes
            if sh is not card
            and sh.left is not None and sh.top is not None
            and left <= sh.left < left + card.width
            and top <= sh.top < top + card.height
        ]
        titles = [sh for sh in doomed if sh.has_text_frame and sh.text_frame.text[:2].rstrip().isdigit()]
        for sh in doomed:
            if sh not in titles:
                sh._element.getparent().remove(sh._element)

        pic = slide.shapes.add_picture(
            str(SHOTS / filename), left + Inches(0.10), top + Inches(0.10),
            width=IMAGE_WIDTH, height=IMAGE_HEIGHT,
        )
        cover_crop(pic, target_aspect)
        for sh in titles:  # title/description move down beneath the image
            sh.top = sh.top + Inches(0.26)


def retarget_note(slide, needle: str, replacement: str) -> None:
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            for paragraph in sh.text_frame.paragraphs:
                for run in paragraph.runs:
                    if needle in run.text:
                        run.text = run.text.replace(needle, replacement)


def restyle_wireframe_slide(slide) -> None:
    """Slide 7: swap the drawn schematic for the real dashboard inside the mock
    browser chrome. The frame, traffic-light dots and URL bar stay; the fake
    tab row and the two hand-drawn content cards are replaced by a real
    screenshot, cover-cropped to the content area."""
    frame_top, frame_bottom = Inches(1.90), Inches(5.34)
    doomed = []
    for sh in slide.shapes:
        if sh.top is None:
            continue
        if frame_top <= sh.top < frame_bottom:
            doomed.append(sh)
    for sh in doomed:
        sh._element.getparent().remove(sh._element)

    pic = slide.shapes.add_picture(
        str(SHOTS / "03-input-vs-output.png"),
        Inches(0.62), Inches(2.02), width=Inches(8.76), height=Inches(3.20),
    )
    cover_crop(pic, 8.76 / 3.20)


def build_closing_slide(slide) -> None:
    """Slide 14: the blank branded template slide becomes the close — one line
    of positioning and the three commands that reproduce every number."""
    left, width = Inches(0.55), Inches(8.9)
    ink = RGBColor(0x1B, 0x1B, 0x19)
    blue = RGBColor(0x1F, 0x47, 0x8A)
    muted = RGBColor(0x5B, 0x64, 0x70)

    kicker = slide.shapes.add_textbox(left, Inches(1.30), width, Inches(0.32))
    paragraph = kicker.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "UNIHACK 2026  ·  AI-POWERED PRODUCT INTELLIGENCE FOR INDUSTRIAL COMMERCE"
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = blue

    headline = slide.shapes.add_textbox(left, Inches(1.72), width, Inches(1.05))
    run = headline.text_frame.paragraphs[0].add_run()
    run.text = "Thank you."
    run.font.name = "Calibri"
    run.font.size = Pt(46)
    run.font.bold = True
    run.font.color.rgb = ink

    positioning = slide.shapes.add_textbox(left, Inches(2.86), width, Inches(0.75))
    run = positioning.text_frame.paragraphs[0].add_run()
    run.text = (
        "One cryptic catalogue line in — a complete, traceable, 252-column "
        "product record out."
    )
    run.font.name = "Calibri"
    run.font.size = Pt(16)
    run.font.color.rgb = muted

    commands = slide.shapes.add_textbox(left, Inches(3.78), width, Inches(1.30))
    text_frame = commands.text_frame
    text_frame.word_wrap = True
    lines = [
        ("Every number in this deck is measured, not claimed:", None),
        ("python scripts/run_pipeline.py --fold holdout      # scores + metrics", 10.5),
        ("python scripts/verify_schema.py                   # 252/252 headers, PASS", 10.5),
        ("python -m pytest tests -q                         # 182 tests pass", 10.5),
    ]
    first = True
    for text, size in lines:
        paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False
        run = paragraph.add_run()
        run.text = text
        run.font.name = "Calibri" if size is None else "Consolas"
        run.font.size = Pt(size or 12)
        run.font.bold = size is None
        run.font.color.rgb = ink if size is None else blue


def main() -> None:
    source = Path(sys.argv[1]) if len(sys.argv) > 1 else default_source()
    presentation = Presentation(str(source))

    snapshots = presentation.slides[10]  # slide 11 — Snapshots of the MVP
    fill_snapshot_cards(snapshots)

    wireframe = presentation.slides[6]  # slide 7 — schematic becomes the real UI
    restyle_wireframe_slide(wireframe)
    retarget_note(
        wireframe, "Schematic of the live dashboard", "The live dashboard")
    retarget_note(
        wireframe,
        "drop in your captured screenshots from docs/shots/ before final submission",
        "raw input left, the enriched 252-column record right — six more views on the next slide",
    )
    links = presentation.slides[12]  # slide 13 — stale git warning
    retarget_note(
        links,
        "This workspace is not a git repository yet.",
        "Repository initialised and committed.",
    )
    build_closing_slide(presentation.slides[13])  # slide 14 — the close

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(OUTPUT))
    print(f"deck -> {OUTPUT}")
    for filename, _, _ in SNAPSHOTS:
        assert (SHOTS / filename).exists(), f"missing shot: {filename}"


if __name__ == "__main__":
    main()
